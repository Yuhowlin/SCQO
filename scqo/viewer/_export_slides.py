"""PowerPoint / Google Slides (.pptx) presentation export builder.

Generates a 16:9 presentation document following the lab's chip presentation guide:
- Slide 1: Title slide (Chip name, Presenter, Cooldown / Setup info)
- Slide 2: Information slide (Design name, Goal, Design pattern)
- Slide 3: Result slide (Predicted vs Measured Qubit Frequency, Design vs Measured Resonator Frequency charts)
- Slide 4: Result slide (Native editable Characterization Summary table)

Compatible with Microsoft PowerPoint and Google Slides (upload to Google Drive
and open as Google Slides for full online editing).
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from ._export_lab_dashboard import _format_mean_std, extract_chip_metrics

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _render_qubit_chart(qubits: list[str], pred_freqs: list[float | None],
                        meas_freqs: list[float | None]) -> bytes:
    """Render the 'Predicted vs Measured Qubit Frequency' chart as PNG bytes."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(6.2, 3.2), dpi=150)
    x = list(range(len(qubits)))

    has_pred = any(v is not None for v in pred_freqs)
    has_meas = any(v is not None for v in meas_freqs)

    if has_pred:
        px = [i for i, v in enumerate(pred_freqs) if v is not None]
        py = [v for v in pred_freqs if v is not None]
        ax.scatter(px, py, color="#2b5c8f", s=40, label="Predicted (R to fq)", zorder=3)
        for xi, yi in zip(px, py):
            ax.annotate(f"{yi:.3f}", (xi, yi), textcoords="offset points",
                        xytext=(0, 6), ha="center", fontsize=7.5, color="#2b5c8f")

    if has_meas:
        mx = [i for i, v in enumerate(meas_freqs) if v is not None]
        my = [v for v in meas_freqs if v is not None]
        ax.scatter(mx, my, color="#d97724", s=45, label="Measured", zorder=4)
        for xi, yi in zip(mx, my):
            ax.annotate(f"{yi:.3f}", (xi, yi), textcoords="offset points",
                        xytext=(0, -12), ha="center", fontsize=7.5, fontweight="bold", color="#d97724")

    ax.set_title("Predicted vs Measured Qubit Frequency", fontsize=10.5, fontweight="bold", pad=8)
    ax.set_ylabel("Frequency (GHz)", fontsize=9)
    ax.set_xticks(x)
    ax.set_xticklabels(qubits, fontsize=8.5)
    ax.grid(True, linestyle="--", alpha=0.4, zorder=1)
    ax.set_axisbelow(True)
    if has_pred or has_meas:
        ax.legend(loc="upper right", fontsize=7.5, framealpha=0.8)

    # Set reasonable y limits
    all_vals = [v for v in pred_freqs + meas_freqs if v is not None]
    if all_vals:
        ymin = max(0.0, min(all_vals) - 0.5)
        ymax = max(all_vals) + 0.6
        ax.set_ylim(ymin, ymax)
    else:
        ax.set_ylim(0, 6)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    return buf.getvalue()


def _render_resonator_chart(qubits: list[str], design_freqs: list[float | None],
                            meas_freqs: list[float | None]) -> bytes:
    """Render the 'Design vs Measured Resonator Frequency' chart as PNG bytes."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(6.2, 3.2), dpi=150)
    x = list(range(len(qubits)))

    has_des = any(v is not None for v in design_freqs)
    has_meas = any(v is not None for v in meas_freqs)

    if has_des:
        dx = [i for i, v in enumerate(design_freqs) if v is not None]
        dy = [v for v in design_freqs if v is not None]
        ax.scatter(dx, dy, color="#2b5c8f", s=40, label="Design", zorder=3)
        for xi, yi in zip(dx, dy):
            ax.annotate(f"{yi:.3f}", (xi, yi), textcoords="offset points",
                        xytext=(0, 6), ha="center", fontsize=7.5, color="#2b5c8f")

    if has_meas:
        mx = [i for i, v in enumerate(meas_freqs) if v is not None]
        my = [v for v in meas_freqs if v is not None]
        ax.scatter(mx, my, color="#d97724", s=45, label="Measured", zorder=4)
        for xi, yi in zip(mx, my):
            ax.annotate(f"{yi:.3f}", (xi, yi), textcoords="offset points",
                        xytext=(0, -12), ha="center", fontsize=7.5, fontweight="bold", color="#d97724")

    ax.set_title("Design vs Measured Resonator Frequency", fontsize=10.5, fontweight="bold", pad=8)
    ax.set_ylabel("Frequency (GHz)", fontsize=9)
    ax.set_xlabel("Resonator", fontsize=9)
    ax.set_xticks(x)
    ax.set_xticklabels(qubits, fontsize=8.5)
    ax.grid(True, linestyle="--", alpha=0.4, zorder=1)
    ax.set_axisbelow(True)
    if has_des or has_meas:
        ax.legend(loc="upper right", fontsize=7.5, framealpha=0.8)

    all_vals = [v for v in design_freqs + meas_freqs if v is not None]
    if all_vals:
        ymin = max(0.0, min(all_vals) - 0.4)
        ymax = max(all_vals) + 0.5
        ax.set_ylim(ymin, ymax)
    else:
        ax.set_ylim(4, 8)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    return buf.getvalue()


def _render_chip_pattern_placeholder(qubits: list[str]) -> bytes:
    """Render a clean schematic diagram for the chip layout pattern."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.patches as patches

    fig, ax = plt.subplots(figsize=(4.5, 4.5), dpi=150)
    ax.set_xlim(-0.1, 1.1)
    ax.set_ylim(-0.1, 1.1)
    ax.axis("off")

    # Outer chip frame
    rect = patches.FancyBboxPatch((0.02, 0.02), 0.96, 0.96,
                                  boxstyle="round,pad=0.02,rounding_size=0.03",
                                  edgecolor="#7d8a99", facecolor="#f8fafc", linewidth=1.5)
    ax.add_patch(rect)

    ax.text(0.5, 0.92, "Design Pattern", ha="center", va="center",
            fontsize=11, fontweight="bold", color="#334155")

    # Grid of qubits
    n = len(qubits)
    cols = min(4, n) if n > 0 else 1
    rows = (n + cols - 1) // cols if cols > 0 else 1

    for idx, q in enumerate(qubits):
        r = idx // cols
        c = idx % cols
        cx = 0.2 + (0.6 * c / max(1, cols - 1)) if cols > 1 else 0.5
        cy = 0.75 - (0.55 * r / max(1, rows - 1)) if rows > 1 else 0.5

        q_box = patches.Circle((cx, cy), 0.065, edgecolor="#2b5c8f", facecolor="#e2e8f0", linewidth=1.2)
        ax.add_patch(q_box)
        ax.text(cx, cy, q.upper(), ha="center", va="center", fontsize=8, fontweight="bold", color="#1e293b")

        # Coupler line to neighbor
        if c < cols - 1 and idx + 1 < n:
            next_cx = 0.2 + (0.6 * (c + 1) / max(1, cols - 1))
            ax.plot([cx + 0.065, next_cx - 0.065], [cy, cy], color="#94a3b8", linestyle="--", linewidth=1)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    return buf.getvalue()


def presentation_pptx_bytes(ctx: dict, store: Any = None, data_root: Path | None = None,
                            presenter: str = "", design_name: str = "",
                            goal: str = "", min_repeats: int = 2,
                            estimator: str = "mean", tags: Any = None) -> bytes:
    """Build a 4-slide PowerPoint (.pptx) presentation document."""
    import pptx
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    metrics = extract_chip_metrics(
        ctx, store=store, data_root=data_root,
        min_repeats=min_repeats, estimator=estimator, tags=tags,
    )
    qubits = metrics["qubits"]
    per_q = metrics["per_qubit"]
    dev = metrics["device"]
    cycle = metrics["cycle"]
    cooldown = metrics["cooldown"]
    setup_name = metrics["setup_name"]

    prs = pptx.Presentation()
    # 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Palette
    color_primary = RGBColor(26, 54, 93)     # #1a365d
    color_muted = RGBColor(100, 116, 139)    # #64748b
    color_bg_card = RGBColor(248, 250, 252)  # #f8fafc
    color_border = RGBColor(226, 232, 240)   # #e2e8f0
    color_hdr_bg = RGBColor(241, 245, 249)   # #f1f5f9
    color_text = RGBColor(30, 41, 59)        # #1e293b

    # Prepare chart images
    pred_fq = [per_q[q]["r_to_fq"] for q in qubits]
    meas_fq = [per_q[q]["f_q_ghz"] for q in qubits]
    qubit_chart_png = _render_qubit_chart(qubits, pred_fq, meas_fq)

    des_res_fq = [per_q[q]["f_bare_ghz"] for q in qubits]
    meas_res_fq = [per_q[q]["f_dress_ghz"] for q in qubits]
    res_chart_png = _render_resonator_chart(qubits, des_res_fq, meas_res_fq)

    # Check for custom design pattern image
    pattern_img_path = None
    if data_root is not None:
        for cand in [data_root / dev / "design_pattern.png",
                     data_root / dev / "chip_layout.png",
                     data_root / dev / "layout.png"]:
            if cand.is_file():
                pattern_img_path = cand
                break

    if pattern_img_path is not None:
        pattern_png = pattern_img_path.read_bytes()
    else:
        pattern_png = _render_chip_pattern_placeholder(qubits)

    # ------------------------------------------------------------- Slide 1: Title
    s1 = prs.slides.add_slide(blank_layout)

    # Decorative top wave / accent bar
    top_shape = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    top_shape.fill.solid()
    top_shape.fill.fore_color.rgb = color_primary
    top_shape.line.color.rgb = color_primary

    # Main Title
    tx_box = s1.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.333), Inches(1.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = dev if dev else "Superconducting QPU"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = color_primary
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    date_str = cycle.get("start", "")
    info_parts = []
    if presenter:
        info_parts.append(f"Presenter: {presenter}")
    if cooldown:
        info_parts.append(f"Cooldown: {cooldown}")
    if setup_name:
        info_parts.append(f"Setup: {setup_name}")
    if date_str:
        info_parts.append(str(date_str))

    p2.text = " · ".join(info_parts) if info_parts else "Characterization Report"
    p2.font.size = Pt(20)
    p2.font.color.rgb = color_muted
    p2.alignment = PP_ALIGN.CENTER

    # ------------------------------------------------------- Slide 2: Information
    s2 = prs.slides.add_slide(blank_layout)

    # Slide Title
    tbox = s2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6.0), Inches(0.8))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Information"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color_primary

    # Left Info Card
    info_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.2), Inches(4.8))
    itf = info_box.text_frame
    itf.word_wrap = True

    p = itf.paragraphs[0]
    p.text = "Design:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color_text

    p = itf.add_paragraph()
    p.text = f"• Design name: {design_name or dev}"
    p.font.size = Pt(18)
    p.font.color.rgb = color_muted
    p.space_after = Pt(24)

    p = itf.add_paragraph()
    p.text = "Goal:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color_text

    p = itf.add_paragraph()
    p.text = f"• {goal or 'Device characterization and parameter baseline verification.'}"
    p.font.size = Pt(18)
    p.font.color.rgb = color_muted

    # Right Image Frame
    s2.shapes.add_picture(io.BytesIO(pattern_png), Inches(7.5), Inches(1.3), width=Inches(5.0))

    # ----------------------------------------------------------- Slide 3: Results (Charts)
    s3 = prs.slides.add_slide(blank_layout)

    tbox = s3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(6.0), Inches(0.7))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Result"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color_primary

    # Left: Two Charts
    s3.shapes.add_picture(io.BytesIO(qubit_chart_png), Inches(0.8), Inches(1.2), width=Inches(6.2))
    s3.shapes.add_picture(io.BytesIO(res_chart_png), Inches(0.8), Inches(4.2), width=Inches(6.2))

    # Right: Design Pattern
    s3.shapes.add_picture(io.BytesIO(pattern_png), Inches(7.5), Inches(1.3), width=Inches(5.0))

    # ------------------------------------------------ Slide 4: Results (Table)
    s4 = prs.slides.add_slide(blank_layout)

    tbox = s4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(6.0), Inches(0.7))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Result"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color_primary

    # Table formatting helper
    def fmt_v(v, unit="", fmt_str="%.3f"):
        if v is None or v == "":
            return "X"
        if isinstance(v, (int, float)):
            return (fmt_str % v) + unit
        return str(v)

    def fmt_stat_or_x(val_str):
        return val_str if val_str else "X"

    # Build candidate table rows with distinct semantics for #100 (statistics) and #1 (single measurement)
    candidate_rows = [
        ("Characterization summary", [q.upper() for q in qubits]),
        ("Dressed resonant freq (GHz)", [fmt_v(per_q[q]["f_dress_ghz"]) for q in qubits]),
        ("Qubit frequency (GHz)", [fmt_v(per_q[q]["f_q_ghz"]) for q in qubits]),
        ("Tunable", ["TRUE" if per_q[q]["tunable"] else "FALSE" for q in qubits]),
        # T1 rows: #100 (statistics) and #1 (single measurement, only if no campaign for this qubit)
        ("T1 #100 (us)", [fmt_stat_or_x(_format_mean_std(per_q[q]["t1_mean_us"], per_q[q]["t1_std_us"], 2)) for q in qubits]),
        ("T1 #1 (us)", [fmt_v(per_q[q]["t1_single_us"] if per_q[q]["t1_mean_us"] is None else None, fmt_str="%.2f") for q in qubits]),
        # Ramsey rows: #100 (statistics) and #1 (single measurement, only if no campaign for this qubit)
        ("T2_Ramsey #100 (us)", [fmt_stat_or_x(_format_mean_std(per_q[q]["t2_ramsey_mean_us"], per_q[q]["t2_ramsey_std_us"], 2)) for q in qubits]),
        ("T2_Ramsey #1 (us)", [fmt_v(per_q[q]["t2_ramsey_single_us"] if per_q[q]["t2_ramsey_mean_us"] is None else None, fmt_str="%.2f") for q in qubits]),
        # Echo rows: #100 (statistics) and #1 (single measurement, only if no campaign for this qubit)
        ("T2_echo #100 (us)", [fmt_stat_or_x(_format_mean_std(per_q[q]["t2_echo_mean_us"], per_q[q]["t2_echo_std_us"], 2)) for q in qubits]),
        ("T2_echo #1 (us)", [fmt_v(per_q[q]["t2_echo_single_us"] if per_q[q]["t2_echo_mean_us"] is None else None, fmt_str="%.2f") for q in qubits]),
        # Coherence ratios
        ("T2*/T1", [fmt_v(per_q[q]["t2_star_over_t1"], fmt_str="%.2f") for q in qubits]),
        ("T2e/T1", [fmt_v(per_q[q]["t2_echo_over_t1"], fmt_str="%.2f") for q in qubits]),
        # Readout fidelity rows: #100 (statistics) and single measurement (only if no campaign for this qubit)
        ("Readout fidelity #100", [fmt_stat_or_x(_format_mean_std(per_q[q]["readout_fidelity_mean"], per_q[q]["readout_fidelity_std"], 2, is_pct=True)) for q in qubits]),
        ("Readout fidelity", [fmt_stat_or_x(_format_mean_std(per_q[q]["readout_fidelity_single"] if per_q[q]["readout_fidelity_mean"] is None else None, None, 2, is_pct=True)) for q in qubits]),
        # Thermal population, temperature & RB fidelity
        ("Temperature #100 (mk)", [fmt_stat_or_x(_format_mean_std(per_q[q]["temperature_mk"], per_q[q]["temperature_err_mk"], 1)) for q in qubits]),
        ("thermal_population #100", [fmt_stat_or_x(_format_mean_std(per_q[q]["thermal_population_mean"], per_q[q]["thermal_population_std"], 2, is_pct=True)) for q in qubits]),
        ("RB fidelity #100", [fmt_stat_or_x(_format_mean_std(per_q[q]["rb_fidelity_mean"], per_q[q]["rb_fidelity_std"], 2, is_pct=True)) for q in qubits]),
        # Physical derivations
        ("Resonator shift Δfr (MHz)", [fmt_v(per_q[q]["delta_fr_mhz"], fmt_str="%.2f") for q in qubits]),
        ("κ/2π (MHz)", [fmt_v(per_q[q]["kappa_mhz"], fmt_str="%.2f") for q in qubits]),
        ("g/2π (MHz, approx.)", [fmt_v(per_q[q]["g_mhz"], fmt_str="%.2f") for q in qubits]),
        ("α/2π (MHz)", [fmt_v(per_q[q]["anharmonicity_mhz"], fmt_str="%.2f") for q in qubits]),
        ("EJ/EC", [fmt_v(per_q[q]["ej_ec"], fmt_str="%.2f") for q in qubits]),
    ]

    # Filter out rows where ALL qubit columns are "X" or empty (keep header row)
    table_rows_spec = []
    for r_title, r_vals in candidate_rows:
        if r_title == "Characterization summary":
            table_rows_spec.append((r_title, r_vals))
        else:
            if any(v not in ("X", "", None) for v in r_vals):
                table_rows_spec.append((r_title, r_vals))

    n_rows = len(table_rows_spec)
    n_cols = 1 + len(qubits)

    tbl_height = min(5.6, 0.42 * n_rows)
    tbl_shape = s4.shapes.add_table(n_rows, n_cols, Inches(0.8), Inches(1.3), Inches(6.5), Inches(tbl_height))
    table = tbl_shape.table

    # Column widths
    table.columns[0].width = Inches(2.5)
    col_w = Inches(4.0 / max(1, len(qubits)))
    for c in range(1, n_cols):
        table.columns[c].width = col_w

    for r_idx, (r_title, r_vals) in enumerate(table_rows_spec):
        # Header / Title column
        c0 = table.cell(r_idx, 0)
        c0.text = r_title
        for p in c0.text_frame.paragraphs:
            p.font.size = Pt(8.5)
            p.font.bold = (r_idx == 0)
            p.font.color.rgb = color_text
            p.alignment = PP_ALIGN.LEFT

        # Data columns
        for c_idx, val in enumerate(r_vals):
            cell = table.cell(r_idx, c_idx + 1)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8.5)
                p.font.bold = (r_idx == 0)
                p.font.color.rgb = color_text
                p.alignment = PP_ALIGN.CENTER

    # Right: Design Pattern
    s4.shapes.add_picture(io.BytesIO(pattern_png), Inches(7.5), Inches(1.3), width=Inches(5.0))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

