"""The 4-slide characterization deck — pure data -> bytes, no fastapi.

python-pptx is lazy-imported (the viewer extra) so its ImportError reaches the
route as the reinstall hint; matplotlib is lazy and forced to Agg, the
``cli/_campaign_plot`` discipline the rest of the viewer already follows.

Two rules the charts hold:

* **A series is labelled what it IS.** The resonator chart's "Design" series is
  the datasheet value from design.toml, never the MEASURED bare frequency —
  plotting the second under the first's name invites the reader to conclude
  that fabrication hit its target when nothing of the sort was checked.
* **Nothing is drawn that was not measured.** A chip layout is either the
  lab's own image or absent; this module will not invent a plausible-looking
  topology and caption it as the design.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from .metrics import extract_chip_metrics
from .xlsx import _format_mean_std

PPTX_MEDIA_TYPE = ("application/vnd.openxmlformats-officedocument"
                   ".presentationml.presentation")

#: filenames a lab may drop in the device folder to supply the real layout.
LAYOUT_FILES = ("design_pattern.png", "chip_layout.png", "layout.png")

_DESIGN_C, _MEASURED_C = "#2b5c8f", "#d97724"


def _pyplot():
    """Headless matplotlib, one place — the three charts shared this preamble
    three times over."""
    import matplotlib

    matplotlib.use("Agg")  # file only: never try to open a window on a lab PC
    import matplotlib.pyplot as plt

    return plt


def _scatter_chart(labels: list[str], series: list[tuple[str, list[float | None], str]],
                   *, title: str, ylabel: str, xlabel: str = "",
                   fallback_ylim: tuple[float, float] = (0.0, 6.0)) -> bytes:
    """One PNG comparing several same-length series over ``labels``.

    Each series is ``(name, values, colour)`` with None for a missing point;
    the x position is the INDEX, so a gap stays a gap rather than shifting the
    points after it onto the wrong qubit.
    """
    plt = _pyplot()
    fig, ax = plt.subplots(figsize=(6.2, 3.2), dpi=150)

    drawn = False
    for offset, (name, values, colour) in enumerate(series):
        pts = [(i, v) for i, v in enumerate(values) if v is not None]
        if not pts:
            continue
        drawn = True
        xs, ys = [p[0] for p in pts], [p[1] for p in pts]
        ax.scatter(xs, ys, color=colour, s=40 + 5 * offset, label=name, zorder=3 + offset)
        for x, y in pts:
            ax.annotate(f"{y:.3f}", (x, y), textcoords="offset points",
                        xytext=(0, 6 if offset == 0 else -12), ha="center",
                        fontsize=7.5, color=colour,
                        fontweight="normal" if offset == 0 else "bold")

    ax.set_title(title, fontsize=10.5, fontweight="bold", pad=8)
    ax.set_ylabel(ylabel, fontsize=9)
    if xlabel:
        ax.set_xlabel(xlabel, fontsize=9)
    ax.set_xticks(list(range(len(labels))))
    ax.set_xticklabels(labels, fontsize=8.5)
    ax.grid(True, linestyle="--", alpha=0.4, zorder=1)
    ax.set_axisbelow(True)
    if drawn:
        ax.legend(loc="upper right", fontsize=7.5, framealpha=0.8)

    values = [v for _n, vals, _c in series for v in vals if v is not None]
    if values:
        ax.set_ylim(max(0.0, min(values) - 0.5), max(values) + 0.6)
    else:
        ax.set_ylim(*fallback_ylim)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    return buf.getvalue()


def find_layout_image(data_root: Path | None, device: str) -> Path | None:
    """The lab's own chip-layout image, or None. There is no generated
    fallback on purpose — see the module docstring."""
    if data_root is None or not device:
        return None
    for name in LAYOUT_FILES:
        candidate = Path(data_root) / device / name
        if candidate.is_file():
            return candidate
    return None


def presentation_pptx_bytes(ctx: dict, store: Any = None, data_root: Path | None = None,
                            presenter: str = "", design_name: str = "", goal: str = "",
                            min_repeats: int = 2, estimator: str = "mean",
                            tags: Any = None, design: Any = None,
                            predicted_f_q: dict[str, float] | None = None) -> bytes:
    import pptx
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    metrics = extract_chip_metrics(
        ctx, store=store, data_root=data_root, min_repeats=min_repeats,
        estimator=estimator, tags=tags, design=design, predicted_f_q=predicted_f_q,
    )
    qubits, per_q = metrics["qubits"], metrics["per_qubit"]
    device, cooldown = metrics["device"], metrics["cooldown"]
    setup_name, cycle = metrics["setup_name"], metrics["cycle"]

    primary = RGBColor(26, 54, 93)
    muted = RGBColor(100, 116, 139)
    text_c = RGBColor(30, 41, 59)

    qubit_png = _scatter_chart(
        qubits,
        [("Predicted (R to fq)", [per_q[q]["r_to_fq"] for q in qubits], _DESIGN_C),
         ("Measured", [per_q[q]["f_q_ghz"] for q in qubits], _MEASURED_C)],
        title="Predicted vs Measured Qubit Frequency", ylabel="Frequency (GHz)")
    # "Design" is the DATASHEET value. per_q[...]["f_bare_ghz"] is measured and
    # belongs to the measured series' family, not this one.
    res_png = _scatter_chart(
        qubits,
        [("Design", [per_q[q]["design_f_bare_ghz"] for q in qubits], _DESIGN_C),
         ("Measured (dressed)", [per_q[q]["f_dress_ghz"] for q in qubits], _MEASURED_C)],
        title="Design vs Measured Resonator Frequency", ylabel="Frequency (GHz)",
        xlabel="Resonator", fallback_ylim=(4.0, 8.0))

    layout_path = find_layout_image(data_root, device)
    layout_png = layout_path.read_bytes() if layout_path else None

    prs = pptx.Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def textbox(slide, left, top, width, height, text, size, *, bold=False,
                color=text_c, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        para.text = text
        para.font.size, para.font.bold, para.font.color.rgb = Pt(size), bold, color
        para.alignment = align
        return frame

    def layout_or_note(slide, left, top, width):
        """The layout image where the lab supplied one; otherwise a line saying
        so, naming the filenames that would be picked up."""
        if layout_png is not None:
            slide.shapes.add_picture(io.BytesIO(layout_png), Inches(left),
                                     Inches(top), width=Inches(width))
        else:
            textbox(slide, left, top, width, 0.9,
                    "No chip layout on file — add "
                    f"{LAYOUT_FILES[0]} to the device folder.",
                    11, color=muted)

    # 1 — title
    s1 = prs.slides.add_slide(blank)
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                              Inches(13.333), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar.line.color.rgb = primary
    frame = textbox(s1, 1.5, 2.8, 10.333, 1.5, device or "Superconducting QPU", 44,
                    bold=True, color=primary, align=PP_ALIGN.CENTER)
    parts = [p for p in (f"Presenter: {presenter}" if presenter else "",
                         f"Cooldown: {cooldown}" if cooldown else "",
                         f"Setup: {setup_name}" if setup_name else "",
                         str(cycle.get("start", "") or "")) if p]
    sub = frame.add_paragraph()
    sub.text = " · ".join(parts) if parts else "Characterization Report"
    sub.font.size, sub.font.color.rgb = Pt(20), muted
    sub.alignment = PP_ALIGN.CENTER

    # 2 — information
    s2 = prs.slides.add_slide(blank)
    textbox(s2, 0.8, 0.5, 6.0, 0.8, "Information", 28, bold=True, color=primary)
    info = textbox(s2, 0.8, 1.5, 6.2, 4.8, "Design:", 22, bold=True)
    for text, size, colour, gap in (
            (f"• Design name: {design_name or device}", 18, muted, 24),
            ("Goal:", 22, text_c, 0),
            (f"• {goal or 'Device characterization and parameter baseline verification.'}",
             18, muted, 0)):
        para = info.add_paragraph()
        para.text, para.font.size, para.font.color.rgb = text, Pt(size), colour
        para.font.bold = size == 22
        if gap:
            para.space_after = Pt(gap)
    layout_or_note(s2, 7.5, 1.3, 5.0)

    # 3 — charts
    s3 = prs.slides.add_slide(blank)
    textbox(s3, 0.8, 0.4, 6.0, 0.7, "Result", 28, bold=True, color=primary)
    s3.shapes.add_picture(io.BytesIO(qubit_png), Inches(0.8), Inches(1.2), width=Inches(6.2))
    s3.shapes.add_picture(io.BytesIO(res_png), Inches(0.8), Inches(4.2), width=Inches(6.2))
    layout_or_note(s3, 7.5, 1.3, 5.0)

    # 4 — summary table
    s4 = prs.slides.add_slide(blank)
    textbox(s4, 0.8, 0.4, 6.0, 0.7, "Result", 28, bold=True, color=primary)

    def val(v, fmt="%.3f"):
        return "X" if v is None else (fmt % v if isinstance(v, (int, float)) else str(v))

    def stat(s):
        return s or "X"

    def only_if_no_campaign(q, single_key, mean_key, fmt="%.2f"):
        return val(per_q[q][single_key] if per_q[q][mean_key] is None else None, fmt)

    candidates: list[tuple[str, list[str]]] = [
        ("Characterization summary", [q.upper() for q in qubits]),
        ("Dressed resonant freq (GHz)", [val(per_q[q]["f_dress_ghz"]) for q in qubits]),
        ("Qubit frequency (GHz)", [val(per_q[q]["f_q_ghz"]) for q in qubits]),
        ("Tunable", ["TRUE" if per_q[q]["tunable"] else "FALSE" for q in qubits]),
        ("T1 #100 (us)", [stat(_format_mean_std(per_q[q]["t1_mean_us"], per_q[q]["t1_std_us"])) for q in qubits]),
        ("T1 #1 (us)", [only_if_no_campaign(q, "t1_single_us", "t1_mean_us") for q in qubits]),
        ("T2_Ramsey #100 (us)", [stat(_format_mean_std(per_q[q]["t2_ramsey_mean_us"], per_q[q]["t2_ramsey_std_us"])) for q in qubits]),
        ("T2_Ramsey #1 (us)", [only_if_no_campaign(q, "t2_ramsey_single_us", "t2_ramsey_mean_us") for q in qubits]),
        ("T2_echo #100 (us)", [stat(_format_mean_std(per_q[q]["t2_echo_mean_us"], per_q[q]["t2_echo_std_us"])) for q in qubits]),
        ("T2_echo #1 (us)", [only_if_no_campaign(q, "t2_echo_single_us", "t2_echo_mean_us") for q in qubits]),
        ("T2*/T1", [val(per_q[q]["t2_star_over_t1"], "%.2f") for q in qubits]),
        ("T2e/T1", [val(per_q[q]["t2_echo_over_t1"], "%.2f") for q in qubits]),
        ("Readout fidelity #100", [stat(_format_mean_std(per_q[q]["readout_fidelity_mean"], per_q[q]["readout_fidelity_std"], 2, True)) for q in qubits]),
        ("Readout fidelity", [stat(_format_mean_std(per_q[q]["readout_fidelity_single"] if per_q[q]["readout_fidelity_mean"] is None else None, None, 2, True)) for q in qubits]),
        ("Temperature #100 (mk)", [stat(_format_mean_std(per_q[q]["temperature_mk"], per_q[q]["temperature_err_mk"], 1)) for q in qubits]),
        ("thermal_population #100", [stat(_format_mean_std(per_q[q]["thermal_population_mean"], per_q[q]["thermal_population_std"], 2, True)) for q in qubits]),
        ("RB fidelity #100", [stat(_format_mean_std(per_q[q]["rb_fidelity_mean"], per_q[q]["rb_fidelity_std"], 2, True)) for q in qubits]),
        ("Resonator shift Δfr (MHz)", [val(per_q[q]["delta_fr_mhz"], "%.2f") for q in qubits]),
        ("κ/2π (MHz)", [val(per_q[q]["kappa_mhz"], "%.2f") for q in qubits]),
        ("g/2π (MHz, approx.)", [val(per_q[q]["g_mhz"], "%.2f") for q in qubits]),
        ("α/2π (MHz)", [val(per_q[q]["anharmonicity_mhz"], "%.2f") for q in qubits]),
        ("EJ/EC", [val(per_q[q]["ej_ec"], "%.2f") for q in qubits]),
    ]
    spec = [(t, v) for t, v in candidates
            if t == "Characterization summary" or any(x not in ("X", "", None) for x in v)]

    table = s4.shapes.add_table(len(spec), 1 + len(qubits), Inches(0.8), Inches(1.3),
                                Inches(6.5), Inches(min(5.6, 0.42 * len(spec)))).table
    table.columns[0].width = Inches(2.5)
    for col in range(1, 1 + len(qubits)):
        table.columns[col].width = Inches(4.0 / max(1, len(qubits)))

    for r, (label, values) in enumerate(spec):
        for c, cell_text in enumerate([label, *values]):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size, para.font.bold = Pt(8.5), (r == 0)
                para.font.color.rgb = text_c
                para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    layout_or_note(s4, 7.5, 1.3, 5.0)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
