"""The lab report's exports, end to end through the viewer.

Reuses ``test_viewer.py``'s ``lab`` fixture — a datastore with applied runs on
two setups of devV's active cycle plus a finished campaign — so these exercise
the real route -> context -> metrics -> renderer path rather than a mock.

The interesting assertions are the ones about what the report must NOT claim:
a missing measurement stays missing, a composite is not a qubit, and a value is
labelled what it is.
"""

from __future__ import annotations

import io

import pytest

from test_viewer import lab  # noqa: F401  (the shared datastore fixture)

from scqo.viewer.lab_report import PPTX_MEDIA_TYPE, XLSX_MEDIA_TYPE
from scqo.viewer.lab_report.metrics import (
    discover_qubits,
    effective_temperature_mk,
    extract_chip_metrics,
)

openpyxl = pytest.importorskip("openpyxl")


def _rows(*entities_fields_values):
    return [{"entity": e, "field": f, "value": v, "unit": "", "source": None,
             "previous": None} for e, f, v in entities_fields_values]


# --------------------------------------------------------------- the routes

def test_setup_dashboard_xlsx_has_the_three_sheets(lab):  # noqa: F811
    resp = lab["client"].get("/setup/devV/cdV/sim_main/export_dashboard.xlsx")
    assert resp.status_code == 200
    assert resp.headers["content-type"] == XLSX_MEDIA_TYPE
    assert "dashboard.xlsx" in resp.headers["content-disposition"]

    book = openpyxl.load_workbook(io.BytesIO(resp.content))
    from scqo.viewer.lab_report.template import load_template

    tpl = load_template()
    assert book.sheetnames == [tpl.sheet("data"), tpl.sheet("dashboard"),
                               tpl.sheet("dictionary")]
    data = book[tpl.sheet("data")]
    assert data["A1"].value == "Basic info"
    assert data["B2"].value == "devV"          # QPU name, from the context
    assert book[tpl.sheet("dictionary")].max_row >= 40


def test_setup_slides_pptx_has_four_slides(lab):  # noqa: F811
    pytest.importorskip("pptx")
    import pptx

    resp = lab["client"].get("/setup/devV/cdV/sim_main/export.pptx")
    assert resp.status_code == 200
    assert resp.headers["content-type"] == PPTX_MEDIA_TYPE

    deck = pptx.Presentation(io.BytesIO(resp.content))
    assert len(deck.slides) == 4
    tables = [s.table for s in deck.slides[3].shapes if s.has_table]
    assert len(tables) == 1
    assert tables[0].cell(0, 0).text == "Characterization summary"


def test_unified_scope_covers_the_whole_cooldown(lab):  # noqa: F811
    for path in ("/cooldown/devV/cdV/export_dashboard.xlsx",
                 "/setup/devV/cdV/sim_main/export_dashboard.xlsx?unified=true"):
        resp = lab["client"].get(path)
        assert resp.status_code == 200, path
        assert "unified" in resp.headers["content-disposition"]


def test_an_unknown_context_is_refused_not_exported(lab):  # noqa: F811
    """The cooldown routes gate their context like the setup page does. An
    export of a chip that does not exist must not come back 200 — it reads as
    'measured nothing' rather than 'no such cooldown'."""
    resp = lab["client"].get("/cooldown/devV/@@@bad@@@/export_dashboard.xlsx")
    assert resp.status_code == 404


def test_a_bad_estimator_is_rejected_not_silently_defaulted(lab):  # noqa: F811
    """`estimator` is typed, so a misspelling is a 422 rather than a workbook
    of mean-based numbers labelled as robust ones."""
    bad = lab["client"].get(
        "/setup/devV/cdV/sim_main/export_dashboard.xlsx?estimator=medain")
    assert bad.status_code == 422
    good = lab["client"].get(
        "/setup/devV/cdV/sim_main/export_dashboard.xlsx?estimator=median")
    assert good.status_code == 200


def test_min_repeats_must_be_positive(lab):  # noqa: F811
    """min_repeats=0 would disable the quality gate while the sheet still says
    '#100'."""
    assert lab["client"].get(
        "/setup/devV/cdV/sim_main/export_dashboard.xlsx?min_repeats=0"
    ).status_code == 422


# --------------------------------------------------------------- the metrics

def test_a_composite_is_not_a_qubit():
    """q1_q2 is a PAIR. Folding it into q1 files two-qubit numbers under one
    qubit — the defect that motivated the catalog-derived suffix set."""
    rows = _rows(("q1", "f_01_hz", 5e9), ("q1_xy", "pi_amp", 0.2),
                 ("q1_res", "f_bare_hz", 7e9), ("q1_q2", "zz_hz", 1e5),
                 ("q2", "f_01_hz", 5.1e9))
    assert discover_qubits(rows) == ["q1", "q2"]


def test_the_sweet_spot_maximum_never_stands_in_for_the_idle_frequency():
    """f_q_max_hz is the arch top; f_01_hz is where the qubit actually sits.
    With only the former, the qubit frequency is REPORTED MISSING — every
    derived quantity depends on it."""
    ctx = {"device": "d", "cooldown": "c", "setup_name": "s", "cycle": {},
           "state_rows": [],
           "physical_rows": _rows(("q1", "f_q_max_hz", 5.5e9),
                                  ("q1", "anharmonicity_hz", -2e8))}
    q1 = extract_chip_metrics(ctx)["per_qubit"]["q1"]
    assert q1["f_q_ghz"] is None
    assert q1["f_q_max_ghz"] == pytest.approx(5.5)
    for derived in ("g_mhz", "ej_ec", "f02_half_ghz", "delta_fr_mhz", "temperature_mk"):
        assert q1[derived] is None, derived


def test_an_empty_context_reports_no_qubits():
    """Not a phantom q1 with every cell blank, which reads as a measured
    failure rather than an empty context."""
    ctx = {"device": "d", "cooldown": "c", "setup_name": "s", "cycle": {},
           "state_rows": [], "physical_rows": []}
    metrics = extract_chip_metrics(ctx)
    assert metrics["qubits"] == [] and metrics["per_qubit"] == {}


def test_a_zero_measurement_is_a_measurement():
    """0.0 is falsy; an `a or b` chain reads it as absent and blanks the cell."""
    ctx = {"device": "d", "cooldown": "c", "setup_name": "s", "cycle": {},
           "state_rows": [], "physical_rows": _rows(("q1", "n_th", 0.0),
                                                    ("q1", "f_01_hz", 5e9))}
    q1 = extract_chip_metrics(ctx)["per_qubit"]["q1"]
    assert q1["thermal_population_mean"] == 0.0
    assert q1["temperature_mk"] is None      # T is undefined at n_th = 0


def test_tunable_reads_the_value_not_the_key():
    """_param_rows emits a row for every observed field, so an unset idle_flux
    still produces a key."""
    base = {"device": "d", "cooldown": "c", "setup_name": "s", "cycle": {},
            "physical_rows": _rows(("q1", "f_01_hz", 5e9))}
    unset = extract_chip_metrics({**base, "state_rows": _rows(("q1_z", "idle_flux", None))})
    assert unset["per_qubit"]["q1"]["tunable"] is False
    setv = extract_chip_metrics({**base, "state_rows": _rows(("q1_z", "idle_flux", 0.1))})
    assert setv["per_qubit"]["q1"]["tunable"] is True


def test_effective_temperature_and_its_error():
    t, err = effective_temperature_mk(0.02, 5.0, 0.004)
    import math

    expected = 47.9924 * 5.0 / math.log1p(1 / 0.02)
    assert t == pytest.approx(expected)
    assert err == pytest.approx(t * 0.004 / (0.02 * 1.02 * math.log1p(1 / 0.02)))
    assert effective_temperature_mk(0.0, 5.0) == (None, None)
    assert effective_temperature_mk(0.02, None) == (None, None)


def test_qubit_order_is_numeric():
    assert discover_qubits(_rows(("q10", "f_01_hz", 1.0), ("q2", "f_01_hz", 1.0),
                                 ("q1", "f_01_hz", 1.0))) == ["q1", "q2", "q10"]
